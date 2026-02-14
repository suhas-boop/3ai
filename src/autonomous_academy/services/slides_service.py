from typing import List, Optional
from io import BytesIO
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from autonomous_academy.schemas import Slide
from duckduckgo_search import DDGS

def _get_image_stream(query: str) -> Optional[BytesIO]:
    """Search for an image and return as a BytesIO stream."""
    try:
        with DDGS() as ddgs:
            # Search for images - getting 1 result
            results = list(ddgs.images(query, max_results=1))
            if results:
                image_url = results[0]['image']
                # Download
                response = requests.get(image_url, timeout=5)
                if response.status_code == 200:
                    return BytesIO(response.content)
    except Exception as e:
        print(f"Failed to get image for '{query}': {e}")
    return None

def generate_pptx(slides: List[Slide], filename: str):
    prs = Presentation()
    
    # Layouts: 0 = Title Slide, 1 = Title and Content
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for i, slide_data in enumerate(slides):
        # Use Title Slide layout for the first slide, Content for the rest
        layout = title_layout if i == 0 else content_layout
        slide = prs.slides.add_slide(layout)
        
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
            
        # ---------------------------------------------------------
        # Title Slide (Index 0)
        # ---------------------------------------------------------
        if i == 0: 
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = "\n".join(slide_data.bullets)

        # ---------------------------------------------------------
        # Content Slides (Index > 0)
        # ---------------------------------------------------------
        else: 
            # 1. Handle Text Content
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                
                # Resize text box to make room for image on the right
                # Standard width is 10 inches. 
                # Let's give text 5.5 inches width (Left side)
                body_shape.left = Inches(0.5)
                body_shape.width = Inches(5.5)
                
                tf = body_shape.text_frame
                tf.clear() # Clear default placeholder text
                tf.word_wrap = True

                for bullet in slide_data.bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(24) # Ensure readable size
                    p.font.name = 'Arial'
                    p.space_after = Pt(10)

            # 2. Handle Image (Right Side)
            if slide_data.image_description:
                 print(f"Searching image for slide: {slide_data.title}...")
                 image_stream = _get_image_stream(slide_data.image_description)
                 if image_stream:
                     try:
                         # Place image on the right
                         # Left = 6.2 inches, Top = 1.5 inches
                         slide.shapes.add_picture(
                             image_stream, 
                             left=Inches(6.2), 
                             top=Inches(1.5), 
                             width=Inches(3.5) # Height auto-scales
                         )
                     except Exception as img_e:
                         print(f"Failed to add image to slide: {img_e}")


        # Add Speaker Notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.speaker_notes

    prs.save(filename)
    return filename

def generate_remote_pptx(topic: str, slide_count: int, filename: str, slides_content: Optional[List[Slide]] = None) -> str:
    """
    Generate a presentation using the external service at localhost:5000.
    
    Args:
        topic: The topic/content prompt for the presentation.
        slide_count: Number of slides to generate.
        filename: Local path to save the downloaded PPTX file.
        slides_content: Optional list of Slide objects to provide detailed context.
        
    Returns:
        The filename of the saved presentation.
    """
    service_url = "http://localhost:5000"
    generate_endpoint = f"{service_url}/api/v1/ppt/presentation/generate"
    
    # Construct a rich prompt if detailed content is available
    prompt_content = f"Make a {slide_count}-slide deck about {topic}."
    
    if slides_content:
        prompt_content += "\n\nUse the following detailed outline and content for the slides:\n"
        for i, slide in enumerate(slides_content):
            prompt_content += f"\nSlide {i+1}: {slide.title}\n"
            if slide.bullets:
                prompt_content += "Key Points:\n" + "\n".join(f"- {b}" for b in slide.bullets) + "\n"
            if slide.speaker_notes:
                prompt_content += f"Speaker Notes: {slide.speaker_notes}\n"
    
    payload = {
        "content": prompt_content,
        "n_slides": slide_count,
        "language": "English",
        "template": "general",
        "export_as": "pptx"
    }
    
    try:
        # 1. Generate
        print(f"Requesting presentation generation for: {topic}")
        # print(f"DEBUG PROMPT: {prompt_content[:500]}...")
        response = requests.post(generate_endpoint, json=payload, timeout=120)
        response.raise_for_status()
        
        data = response.json()
        remote_path = data.get("path")
        
        if not remote_path:
            raise ValueError(f"No path returned from presentation service: {data}")
            
        # 2. Download
        # The service returns a path like /app_data/exports/Filename.pptx
        # We need to URL encode it or simply append it if it's a valid path segment
        # In the reproduction script, we found we could access it directly via localhost:5000{path}
        
        # Ensure path starts with /
        if not remote_path.startswith("/"):
            remote_path = f"/{remote_path}"
            
        download_url = f"{service_url}{remote_path}"
        print(f"Downloading presentation from: {download_url}")
        
        download_response = requests.get(download_url, timeout=60)
        download_response.raise_for_status()
        
        with open(filename, "wb") as f:
            f.write(download_response.content)
            
        print(f"Presentation saved to: {filename}")
        return filename
        
    except Exception as e:
        print(f"Failed to generate remote presentation: {e}")
        # Build a simple fallback or re-raise depending on requirements.
        # For now, Re-raise so the API knows it failed.
        raise e

